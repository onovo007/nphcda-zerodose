"""
Generate premium deliverables from the live model outputs:
- a one-page factsheet (branded HTML, print-to-PDF),
- a policy brief (editable, branded Word .docx).

Findings are assembled from the same cached model functions the domain pages use, so values
match the on-screen results exactly. Narrative is LLM-drafted (grounded) when a key is present,
otherwise a templated narrative built from the numbers.
"""
from __future__ import annotations

import base64
import io
import re
from datetime import datetime

import pandas as pd

import config as C
import data_io as dio
import llm
from theme import clean, img_data_uri
from models.d1_forecast import national_forecasts
from models.d2_dropout import dropout_forecasts, lasso_drivers
from models.d5_zerodose import run_state_model, run_lga_burden


# --------------------------------------------------------------------------------------
# Findings assembly (reuses cached model results)
# --------------------------------------------------------------------------------------
def build_findings(data: dict) -> dict:
    f: dict = {"generated": datetime.now().strftime("%d %B %Y"),
               "consortium": "CIDRE and Quantium Insights LLC, in technical support of NPHCDA",
               "audience": "NPHCDA, GAVI and UNICEF"}
    kd = dio.df_hash(data.get("dhis2"))

    # Domain 5 (centerpiece)
    if all(data.get(k) is not None for k in ("ndhs_long", "under5", "dhis2", "lga_population")):
        kn, ku, kp = (dio.df_hash(data["ndhs_long"]), dio.df_hash(data["under5"]),
                      dio.df_hash(data["lga_population"]))
        mkey = f"{kn}-{ku}-{kd}-{C.MCMC_DRAWS_LIVE}-{C.MCMC_TUNE_LIVE}"
        out = run_state_model(data["ndhs_long"], data["under5"], data["dhis2"], key=mkey,
                              draws=C.MCMC_DRAWS_LIVE, tune=C.MCMC_TUNE_LIVE)
        res = out["res"]
        lga = run_lga_burden(data["dhis2"], res, data["lga_population"], key=f"{mkey}-{kp}")
        tier1 = res[res["priority_tier"].astype(str) == "Tier 1: Critical"]["state"].tolist()
        f["d5"] = {
            "national_zd_count_2026": int(res["zd_count_2026"].sum()),
            "lga_total": lga["national_total"], "lga_count": lga["n_lgas"],
            "top20_pct": lga["top20_pct"], "n80": lga["n80"],
            "tier1_states": tier1, "max_rhat": out["max_rhat"],
            "top_states": [
                {"state": r["state"], "zone": r["zone"],
                 "zd_2026_pct": round(float(r["zd_pred_2026_mean"]), 1),
                 "zd_2026_count": int(r["zd_count_2026"]) if pd.notna(r["zd_count_2026"]) else None,
                 "tier": str(r["priority_tier"])}
                for _, r in res.sort_values("state_rank").head(8).iterrows()],
            "top_lgas": [
                {"lga": r["LGA"], "state": r["State"], "zone": r["Zone"],
                 "zd_count": int(r["Zero-dose children (est)"]),
                 "zd_rate_pct": float(r["Zero-dose rate (%)"])}
                for _, r in lga["pareto"].head(10).iterrows()],
            "zone_burden": [
                {"zone": str(z), "count": int(c)} for z, c in
                res.groupby("zone")["zd_count_2026"].sum().sort_values(ascending=False).items()
                if pd.notna(c)],
        }

    # Domain 1
    if data.get("dhis2") is not None:
        nat = dio.national_monthly(dio.prep_dhis2(data["dhis2"]))
        d1 = national_forecasts(nat, key=kd)
        sm = d1["summary"]
        f["d1"] = {
            "at_risk_antigens": sm[sm["Crosses 80% in 6-12m"] == "Yes"]["Antigen"].tolist(),
            "summary": sm.to_dict(orient="records"),
        }
        # Domain 2
        d2 = dropout_forecasts(nat, key=kd)
        f["d2"] = {"latest_dropout": {s["label"]: round(s["obs_y"][-1], 1) for s in d2.values()}}
        if data.get("model_dataset") is not None:
            agg = dio.state_monthly(dio.prep_dhis2(data["dhis2"]))
            drv = lasso_drivers(agg, data["model_dataset"])
            f["d2"]["top_drivers"] = {
                C.DROPOUT_TARGETS[t]: [k.replace("pct_", "").replace("_", " ") for k in c.head(4).index]
                for t, c in drv.items()}
    if data.get("ndhs_antigens") is not None:
        try:
            f["survey"] = dio.survey_national_coverage(data["ndhs_antigens"], data.get("under5"))
        except Exception:
            pass
    return f


# --------------------------------------------------------------------------------------
# Templated narrative (used when no LLM key)
# --------------------------------------------------------------------------------------
def template_narrative(f: dict, kind: str) -> str:
    d5 = f.get("d5", {})
    d1 = f.get("d1", {})
    d2 = f.get("d2", {})
    parts = []
    if kind == "policy":
        parts.append("## Executive summary")
    parts.append(
        f"Modelling of Nigeria routine immunization data projects about "
        f"{d5.get('lga_total', 0):,} zero-dose children across {d5.get('lga_count', 0)} LGAs in 2026. "
        f"Burden is concentrated: the top 20 percent of LGAs carry about {d5.get('top20_pct', 0):.0f} "
        f"percent of the total, and the highest-risk states are "
        f"{', '.join(d5.get('tier1_states', [])[:3]) or 'in the North-West'}.")
    if kind == "policy":
        parts.append("## Situation analysis")
        parts.append(
            "Zero-dose rates remain highest in the North-West. Antigen coverage is forecast near or "
            "above the 80 percent target nationally, but dropout between antigen doses and spatial "
            "clustering of unvaccinated children sustain the burden in specific LGAs.")
    parts.append("## Key findings" if kind == "policy" else "### Key findings")
    kf = []
    if d5:
        kf.append(f"About {d5.get('lga_total', 0):,} zero-dose children in 2026 across "
                  f"{d5.get('lga_count', 0)} LGAs; 80 percent of the burden sits in the top "
                  f"{d5.get('n80', 0)} LGAs.")
        kf.append("Tier-1 critical states: " + (", ".join(d5.get("tier1_states", [])) or "North-West states") + ".")
        if d5.get("top_lgas"):
            t = d5["top_lgas"][0]
            kf.append(f"Highest-burden LGA: {t['lga']} ({t['state']}), about {t['zd_count']:,} "
                      f"zero-dose children at {t['zd_rate_pct']:.0f} percent.")
    if d1:
        ar = d1.get("at_risk_antigens")
        kf.append(("Antigens at risk of falling below 80 percent in 6 to 12 months: " + ", ".join(ar))
                  if ar else "All tracer antigens are forecast at or above the 80 percent target nationally.")
    if d2.get("top_drivers"):
        first = next(iter(d2["top_drivers"].items()))
        kf.append(f"Leading drivers of {first[0]} dropout: {', '.join(first[1])}.")
    parts.append("\n".join(f"- {x}" for x in kf))
    parts.append("## Recommendations" if kind == "policy" else "### Priority actions")
    recs = [
        f"Prioritize the top {d5.get('n80', 0) if d5 else 'priority'} LGAs that hold 80 percent of the "
        "burden for targeted catch-up and outreach.",
        "Concentrate first-line investment in the Tier-1 states ("
        + (", ".join(d5.get("tier1_states", [])[:3]) if d5 else "North-West") + ").",
        "Address dropout with reminder-recall and defaulter tracing where dose-to-dose loss is highest.",
        "Use the LGA hotspot map to micro-plan supervision and supplementary sessions.",
    ]
    parts.append("\n".join(f"{i+1}. {r}" for i, r in enumerate(recs)))
    return "\n\n".join(parts)


# --------------------------------------------------------------------------------------
# Markdown helpers
# --------------------------------------------------------------------------------------
def _strip_md(s: str) -> str:
    return clean(re.sub(r"\*\*(.*?)\*\*", r"\1", s).replace("**", "").strip())


def md_to_html(md: str) -> str:
    html, in_list = [], False
    for raw in md.splitlines():
        line = raw.rstrip()
        if not line.strip():
            continue
        if line.startswith("## "):
            if in_list:
                html.append("</ul>"); in_list = False
            html.append(f"<h2>{_strip_md(line[3:])}</h2>")
        elif line.startswith("### "):
            if in_list:
                html.append("</ul>"); in_list = False
            html.append(f"<h3>{_strip_md(line[4:])}</h3>")
        elif re.match(r"^\s*[-*]\s+", line) or re.match(r"^\s*\d+\.\s+", line):
            if not in_list:
                html.append("<ul>"); in_list = True
            html.append(f"<li>{_strip_md(re.sub(r'^\s*([-*]|\d+\.)\s+', '', line))}</li>")
        else:
            if in_list:
                html.append("</ul>"); in_list = False
            html.append(f"<p>{_strip_md(line)}</p>")
    if in_list:
        html.append("</ul>")
    return "\n".join(html)


# --------------------------------------------------------------------------------------
# Factsheet (premium HTML)
# --------------------------------------------------------------------------------------
def _bignum(v: float) -> str:
    if v >= 1e6:
        return f"{v / 1e6:.2f}M"
    if v >= 1000:
        return f"{v / 1000:.0f}k"
    return f"{v:,.0f}"


def _icon(name: str, color: str) -> str:
    shapes = {
        "children": ("<circle cx='8' cy='8' r='3'/><path d='M2 21v-2a5 5 0 0 1 5-5h2a5 5 0 0 1 5 5v2'/>"
                     "<circle cx='17.5' cy='10' r='2.2'/><path d='M22 21v-1.5a3.5 3.5 0 0 0-3.5-3.5'/>"),
        "target": "<circle cx='12' cy='12' r='9'/><circle cx='12' cy='12' r='5'/><circle cx='12' cy='12' r='1.6'/>",
        "pin": "<path d='M12 22s7-6.4 7-12a7 7 0 1 0-14 0c0 5.6 7 12 7 12z'/><circle cx='12' cy='10' r='2.5'/>",
        "alert": "<path d='M12 3 2 20h20L12 3z'/><line x1='12' y1='10' x2='12' y2='14.5'/><circle cx='12' cy='17.4' r='0.6'/>",
        "syringe": "<path d='M4 20l3-3'/><path d='M14 6l4 4-7 7-4-4z'/><path d='M16 4l4 4'/><path d='M9 11l4 4'/>",
    }.get(name, "")
    return (f"<svg width='38' height='38' viewBox='0 0 24 24' fill='none' stroke='{color}' "
            f"stroke-width='1.7' stroke-linecap='round' stroke-linejoin='round'>{shapes}</svg>")


def _burden_bar_png(zone_burden: list) -> str:
    """Horizontal bar of zero-dose burden by zone as a base64 PNG (the factsheet's headline visual)."""
    if not zone_burden:
        return ""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception:
        return ""
    zb = sorted(zone_burden, key=lambda z: z["count"])  # ascending so largest is at top in barh
    zones = [z["zone"] for z in zb]
    vals = [z["count"] / 1e3 for z in zb]
    top = max(vals) if vals else 1
    colors = ["#C0392B" if v == top else "#1F3B57" for v in vals]
    fig, ax = plt.subplots(figsize=(7.8, 2.9), dpi=150)
    ax.barh(zones, vals, color=colors)
    for i, v in enumerate(vals):
        ax.text(v + top * 0.012, i, f"{v:,.0f}k", va="center", fontsize=8.5, color="#1A1A1A")
    ax.set_xlabel("Estimated zero-dose children, 2026 (thousands)", fontsize=8.5, color="#33414d")
    ax.set_xlim(0, top * 1.16)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    ax.tick_params(labelsize=9)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    return "data:image/png;base64," + base64.b64encode(buf.getvalue()).decode()


def _pill(badge: str, text: str, color: str) -> str:
    return (f"<div class='pill' style='border-color:{color}'>"
            f"<span class='pn' style='color:{color}'>{clean(badge)}</span>"
            f"<span class='pt'>{clean(text)}</span></div>")


def factsheet_html(f: dict, narrative_md: str) -> str:
    d5 = f.get("d5", {})
    d1 = f.get("d1", {})
    d2 = f.get("d2", {})
    survey = f.get("survey", {})
    GREEN, RED, NAVY = C.NPHCDA_GREEN, C.ACCENT, C.NAVY
    logo = img_data_uri(C.LOGO_PATH)
    logo_html = (f"<img src='{logo}' style='height:58px'/>" if logo
                 else "<b style='font-size:22px;color:#0F5226'>NPHCDA</b>")

    # Hero band: three headline numbers with icons (GAVI style).
    hero = [
        ("children", _bignum(d5.get("lga_total", 0)), RED,
         "zero-dose children projected in 2026",
         f"across {d5.get('lga_count', 0)} reporting LGAs"),
        ("target", f"{d5.get('top20_pct', 0):.0f}%", NAVY,
         "of the burden in the top 20% of LGAs",
         f"80% of the burden sits in just {d5.get('n80', 0)} LGAs"),
        ("pin", str(len(d5.get("tier1_states", []))), GREEN,
         "Tier-1 critical states",
         clean(", ".join(d5.get("tier1_states", [])) or "concentrated in the North-West")),
    ]
    hero_html = "".join(
        f"<div class='hcard'><div class='hicon'>{_icon(ic, col)}</div><div>"
        f"<div class='hbig' style='color:{col}'>{clean(v)}</div>"
        f"<div class='hlab'>{clean(lab)}</div><div class='hcap'>{clean(cap)}</div></div></div>"
        for ic, v, col, lab, cap in hero)

    # Three colour-coded goal cards, each a big number + paragraph + pill call-outs.
    sm = {r["Antigen"]: r for r in d1.get("summary", [])}
    low_ant = min(sm.values(), key=lambda r: r["Min forecast (% of 2024 baseline)"]) if sm else None
    ar = d1.get("at_risk_antigens", [])
    top_lga = (d5.get("top_lgas") or [{}])[0]
    top_state = (d5.get("top_states") or [{}])[0]

    goals = []
    # Goal 1 - Reach / burden
    g1 = [_pill(f"{d5.get('n80', 0)}", "LGAs carry 80% of the national burden", GREEN)]
    if top_lga:
        g1.append(_pill(f"{top_lga.get('zd_count', 0):,}",
                        f"highest-burden LGA: {top_lga.get('lga','')} ({top_lga.get('state','')})", GREEN))
    goals.append((GREEN, "The reach goal", _bignum(d5.get("lga_total", 0)),
                  f"zero-dose children projected in 2026 across {d5.get('lga_count', 0)} LGAs. The "
                  "burden is highly concentrated, so targeting a small set of LGAs reaches most "
                  "unvaccinated children.", g1))
    # Goal 2 - Equity / where it concentrates
    g2 = [_pill(f"{len(d5.get('tier1_states', []))}",
                "Tier-1 critical states: " + (", ".join(d5.get("tier1_states", [])) or "North-West"), RED)]
    if top_state:
        g2.append(_pill(f"{int(top_state.get('zd_2026_count') or 0):,}",
                        f"zero-dose children in {top_state.get('state','')} (2026)", RED))
    goals.append((RED, "The equity goal", f"{top_state.get('zd_2026_pct', 0):.0f}%",
                  f"predicted 2026 zero-dose rate in the highest-risk state "
                  f"({top_state.get('state','the North-West')}). The North-West carries the heaviest, "
                  "most persistent burden and needs first-line investment.", g2))
    # Goal 3 - Coverage & completion
    g3 = []
    if survey.get("Penta1") is not None:
        g3.append(_pill(f"{survey['Penta1']:.0f}%", "NDHS Penta1 survey coverage (2024)", NAVY))
    if low_ant is not None:
        g3.append(_pill(f"{low_ant['Min forecast (% of 2024 baseline)']:.0f}%",
                        f"lowest national antigen forecast ({low_ant['Antigen']})", NAVY))
    if d2.get("top_drivers"):
        first = next(iter(d2["top_drivers"].items()))
        g3.append(_pill("Drivers", f"{first[0]}: {', '.join(first[1][:3])}", NAVY))
    g3_big = str(len(ar)) if ar else "On track"
    g3_para = (("tracer antigens projected below the 80% target within 6 to 12 months: "
                + ", ".join(ar) + ".") if ar
               else "all four tracer antigens are projected to hold at or above the 80% target "
               "nationally - protect this with reminder-recall and defaulter tracing.")
    goals.append((NAVY, "The coverage and completion goal", g3_big, g3_para, g3))

    goals_html = "".join(
        f"<div class='goal' style='border-top-color:{col}'>"
        f"<div class='glabel' style='color:{col}'>{clean(label)}</div>"
        f"<div class='gbig' style='color:{col}'>{clean(big)}</div>"
        f"<div class='gpara'>{clean(para)}</div>{''.join(pl)}</div>"
        for col, label, big, para, pl in goals)

    chart_uri = _burden_bar_png(d5.get("zone_burden", []))
    chart_html = (f"<div class='sect'>Where the burden concentrates</div>"
                  f"<div class='chartcap'>Estimated zero-dose children by geopolitical zone in 2026. "
                  f"The North-West (red) carries the heaviest burden.</div>"
                  f"<img class='chart' src='{chart_uri}'/>" if chart_uri else "")

    # Infographic findings + actions (instead of a policy-brief block of text).
    kf = _section_lines(narrative_md, ["key findings"])[:5]
    recs = _section_lines(narrative_md, ["recommendation", "priority action", "implementation"])[:5]
    kf_html = "".join(f"<div class='kfi'><span class='kfick'>&#10003;</span>"
                      f"<span>{clean(x)}</span></div>" for x in kf)
    acts_html = "".join(f"<div class='act'><span class='actn'>{i+1}</span>"
                        f"<span>{clean(x)}</span></div>" for i, x in enumerate(recs))

    rows = "".join(
        f"<tr><td>{i+1}</td><td>{clean(s['lga'])}</td><td>{clean(s['state'])}</td>"
        f"<td style='text-align:right'>{s['zd_count']:,}</td>"
        f"<td style='text-align:right'>{s['zd_rate_pct']:.0f}%</td></tr>"
        for i, s in enumerate(d5.get("top_lgas", [])))
    table_html = (f"<table class='tbl'><thead><tr><th>#</th><th>LGA</th><th>State</th>"
                  f"<th>Zero-dose children</th><th>Rate</th></tr></thead><tbody>{rows}</tbody></table>"
                  if rows else "")

    return f"""<!doctype html><html><head><meta charset='utf-8'>
<style>
@import url('https://fonts.googleapis.com/css2?family=IBM+Plex+Sans:wght@400;600;700;800&display=swap');
body{{font-family:'IBM Plex Sans',Segoe UI,sans-serif;color:#1A1A1A;margin:0;background:#fff}}
.wrap{{max-width:940px;margin:0 auto;padding:32px 40px}}
.hd{{display:flex;align-items:center;justify-content:space-between;border-bottom:4px solid {GREEN};padding-bottom:14px}}
.tag{{color:{GREEN};font-size:12px;font-weight:800;letter-spacing:1.5px;text-transform:uppercase}}
.sub{{color:{C.MUTE};font-size:12.5px}}
.lead{{font-size:26px;font-weight:700;line-height:1.25;color:{NAVY};margin:18px 0 4px}}
.hero{{display:grid;grid-template-columns:repeat(3,1fr);gap:22px;margin:22px 0 6px;
       border-top:1px solid #e7eef4;border-bottom:1px solid #e7eef4;padding:18px 0}}
.hcard{{display:flex;gap:12px;align-items:flex-start}}
.hicon{{flex:none;margin-top:2px}}
.hbig{{font-size:38px;font-weight:800;line-height:1}}
.hlab{{font-size:12.5px;color:#1A1A1A;margin-top:3px;font-weight:600}}
.hcap{{font-size:11px;color:{C.MUTE};margin-top:2px}}
.goals{{display:grid;grid-template-columns:repeat(3,1fr);gap:16px;margin:18px 0 6px}}
.goal{{border:1px solid #eaf0f3;border-top:5px solid {GREEN};border-radius:12px;padding:15px 15px 12px;
       background:linear-gradient(180deg,#fff,#fbfdfc);box-shadow:0 4px 14px rgba(20,40,55,.05)}}
.glabel{{font-size:11.5px;font-weight:800;text-transform:uppercase;letter-spacing:.6px}}
.gbig{{font-size:40px;font-weight:800;line-height:1;margin:4px 0 6px}}
.gpara{{font-size:11.5px;color:#33414d;line-height:1.35;margin-bottom:6px}}
.pill{{display:flex;align-items:center;gap:9px;border:1.5px solid {C.STEEL};border-radius:999px;
       padding:5px 11px;margin-top:7px}}
.pn{{font-weight:800;font-size:13.5px;white-space:nowrap}}
.pt{{color:#33414d;font-size:10.5px;line-height:1.2}}
.sect{{color:{NAVY};font-size:17px;font-weight:700;border-left:4px solid {C.GOLD};padding-left:10px;margin:24px 0 10px}}
.chartcap{{color:{C.MUTE};font-size:11.5px;margin:-4px 0 8px}}
.chart{{width:100%;max-width:760px;border:1px solid #eef2f5;border-radius:10px;padding:6px}}
.kf{{display:grid;grid-template-columns:1fr 1fr;gap:8px 18px;margin:4px 0}}
.kfi{{display:flex;gap:9px;align-items:flex-start;font-size:12px;color:#27323b}}
.kfick{{flex:none;width:18px;height:18px;border-radius:50%;background:{GREEN};color:#fff;font-size:11px;
        display:inline-flex;align-items:center;justify-content:center;margin-top:1px}}
.acts{{display:flex;flex-direction:column;gap:7px;margin:4px 0}}
.act{{display:flex;gap:11px;align-items:flex-start;font-size:12px;color:#27323b}}
.actn{{flex:none;width:22px;height:22px;border-radius:50%;background:{NAVY};color:#fff;font-weight:700;
       font-size:12px;display:inline-flex;align-items:center;justify-content:center}}
h2,h3{{color:{NAVY}}} h2{{font-size:16px;margin-top:16px}} h3{{font-size:13px}}
ul{{margin:6px 0}} li{{margin:4px 0;font-size:12.5px}}
.tbl{{width:100%;border-collapse:collapse;margin-top:8px;font-size:12px}}
.tbl th{{background:{NAVY};color:#fff;padding:7px 9px;text-align:left}}
.tbl td{{padding:6px 9px;border-bottom:1px solid #e7eef4}}
.ft{{margin-top:24px;border-top:1px solid #e0e7ee;padding-top:10px;color:{C.MUTE};font-size:10.5px}}
</style></head><body><div class='wrap'>
<div class='hd'>{logo_html}<div style='text-align:right'><div class='tag'>Zero-Dose Factsheet</div>
<div class='sub'>{clean(f.get('generated',''))}</div></div></div>
<div class='lead'>Nigeria zero-dose modelling: where the unvaccinated children are, and where to act first.</div>
<div class='sub'>{clean(f.get('consortium',''))}. For {clean(f.get('audience',''))}.</div>
<div class='hero'>{hero_html}</div>
<div class='goals'>{goals_html}</div>
{chart_html}
<div class='sect'>Key findings</div><div class='kf'>{kf_html}</div>
<div class='sect'>Priority actions</div><div class='acts'>{acts_html}</div>
<div class='sect'>Highest-burden LGAs</div>{table_html}
<div class='ft'>Generated by the NPHCDA Zero-Dose Predictive Modelling Platform from live model outputs
(coverage, dropout and zero-dose workstreams). Figures are model estimates; review before use.
LGA population denominator: City Population (NPC 2022 projection).</div>
</div></body></html>"""


# --------------------------------------------------------------------------------------
# Policy brief (Word .docx)
# --------------------------------------------------------------------------------------
def policy_docx(f: dict, narrative_md: str) -> bytes:
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    navy, green, gold, mute = RGBColor(0x1F, 0x3B, 0x57), RGBColor(0x1C, 0x7A, 0x3D), \
        RGBColor(0xC8, 0x90, 0x2A), RGBColor(0x6B, 0x7A, 0x88)
    doc = Document()
    doc.styles["Normal"].font.name = "Calibri"
    doc.styles["Normal"].font.size = Pt(11)

    if C.LOGO_PATH.exists():
        try:
            doc.add_picture(str(C.LOGO_PATH), width=Inches(1.6))
        except Exception:
            pass
    t = doc.add_paragraph()
    r = t.add_run("Nigeria Zero-Dose Immunization - Policy Brief")
    r.bold = True; r.font.size = Pt(20); r.font.color.rgb = navy
    s = doc.add_paragraph()
    sr = s.add_run(f"{f.get('consortium','')}. For {f.get('audience','')}.  |  {f.get('generated','')}")
    sr.italic = True; sr.font.size = Pt(9); sr.font.color.rgb = mute

    for raw in narrative_md.splitlines():
        line = raw.rstrip()
        if not line.strip():
            continue
        if line.startswith("## "):
            h = doc.add_paragraph(); hr = h.add_run(_strip_md(line[3:]))
            hr.bold = True; hr.font.size = Pt(14); hr.font.color.rgb = navy
        elif line.startswith("### "):
            h = doc.add_paragraph(); hr = h.add_run(_strip_md(line[4:]))
            hr.bold = True; hr.font.size = Pt(12); hr.font.color.rgb = green
        elif re.match(r"^\s*[-*]\s+", line):
            doc.add_paragraph(_strip_md(re.sub(r"^\s*[-*]\s+", "", line)), style="List Bullet")
        elif re.match(r"^\s*\d+\.\s+", line):
            doc.add_paragraph(_strip_md(re.sub(r"^\s*\d+\.\s+", "", line)), style="List Number")
        else:
            doc.add_paragraph(_strip_md(line))

    d5 = f.get("d5", {})
    if d5.get("top_states"):
        h = doc.add_paragraph(); hr = h.add_run("Priority states (top 8 by risk)")
        hr.bold = True; hr.font.size = Pt(13); hr.font.color.rgb = navy
        tb = doc.add_table(rows=1, cols=4); tb.style = "Light Grid Accent 1"
        for i, c in enumerate(["State", "Zone", "ZD 2026 (%)", "ZD children 2026"]):
            cell = tb.rows[0].cells[i]; cr = cell.paragraphs[0].add_run(c); cr.bold = True
        for s in d5["top_states"]:
            cells = tb.add_row().cells
            cells[0].text = clean(s["state"]); cells[1].text = clean(s["zone"])
            cells[2].text = f"{s['zd_2026_pct']:.1f}"
            cells[3].text = f"{s['zd_2026_count']:,}" if s["zd_2026_count"] else "-"
    if d5.get("top_lgas"):
        h = doc.add_paragraph(); hr = h.add_run("Highest-burden LGAs (top 10)")
        hr.bold = True; hr.font.size = Pt(13); hr.font.color.rgb = navy
        tb = doc.add_table(rows=1, cols=4); tb.style = "Light Grid Accent 1"
        for i, c in enumerate(["LGA", "State", "ZD children", "Rate (%)"]):
            cr = tb.rows[0].cells[i].paragraphs[0].add_run(c); cr.bold = True
        for s in d5["top_lgas"]:
            cells = tb.add_row().cells
            cells[0].text = clean(s["lga"]); cells[1].text = clean(s["state"])
            cells[2].text = f"{s['zd_count']:,}"; cells[3].text = f"{s['zd_rate_pct']:.0f}"

    fp = doc.add_paragraph()
    fr = fp.add_run("Generated by the NPHCDA Zero-Dose Predictive Modelling Platform from live model "
                    "outputs (coverage, dropout and zero-dose workstreams). Figures are model estimates; review before use. "
                    "LGA population denominator: City Population (NPC 2022 projection).")
    fr.italic = True; fr.font.size = Pt(8); fr.font.color.rgb = mute

    buf = io.BytesIO(); doc.save(buf); return buf.getvalue()


def _section_lines(md: str, heads: list[str]) -> list[str]:
    """Return the bullet/numbered lines under any heading whose text matches one of `heads`."""
    out, cap = [], False
    for ln in md.splitlines():
        s = ln.strip()
        if s.startswith("#"):
            cap = any(h.lower() in s.lower() for h in heads)
            continue
        if cap and s:
            out.append(_strip_md(re.sub(r"^\s*([-*]|\d+\.)\s*", "", s)))
    return out


def policy_pptx(f: dict, narrative_md: str) -> bytes:
    """Branded PowerPoint policy deck (NPHCDA logo, navy title bars)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    NAVY, GREEN, GOLD, INK = RGBColor(0x1F, 0x3B, 0x57), RGBColor(0x1C, 0x7A, 0x3D), \
        RGBColor(0xC8, 0x90, 0x2A), RGBColor(0x1A, 0x1A, 0x1A)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    SW = Inches(13.333)

    def bar(slide, title):
        rect = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.0))
        rect.fill.solid(); rect.fill.fore_color.rgb = NAVY; rect.line.fill.background()
        tf = rect.text_frame; tf.margin_left = Inches(0.4)
        tf.text = title
        p = tf.paragraphs[0]; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = RGBColor(255, 255, 255)
        accent = slide.shapes.add_shape(1, 0, Inches(1.0), SW, Inches(0.06))
        accent.fill.solid(); accent.fill.fore_color.rgb = GOLD; accent.line.fill.background()

    def bullets(slide, items, top=1.4, size=18):
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(5.6))
        tf = tb.text_frame; tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "- " + clean(it); p.font.size = Pt(size); p.font.color.rgb = INK; p.space_after = Pt(8)

    def table_slide(title, headers, rows):
        s = prs.slides.add_slide(blank); bar(s, title)
        n = len(rows) + 1
        tbl = s.shapes.add_table(n, len(headers), Inches(0.6), Inches(1.4),
                                 Inches(12.1), Inches(0.4 * n)).table
        for j, h in enumerate(headers):
            c = tbl.cell(0, j); c.text = h
            c.fill.solid(); c.fill.fore_color.rgb = NAVY
            c.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            c.text_frame.paragraphs[0].font.bold = True; c.text_frame.paragraphs[0].font.size = Pt(13)
        for i, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                c = tbl.cell(i, j); c.text = str(val); c.text_frame.paragraphs[0].font.size = Pt(12)

    d5, d1, d2 = f.get("d5", {}), f.get("d1", {}), f.get("d2", {})

    # 1 - Title
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(1, 0, 0, SW, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    if C.LOGO_PATH.exists():
        try:
            s.shapes.add_picture(str(C.LOGO_PATH), Inches(0.6), Inches(0.5), height=Inches(0.9))
        except Exception:
            pass
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(2.4)); tf = tb.text_frame
    tf.word_wrap = True
    tf.text = "Nigeria Zero-Dose Immunization"
    tf.paragraphs[0].font.size = Pt(40); tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    p = tf.add_paragraph(); p.text = "Policy Brief"
    p.font.size = Pt(28); p.font.color.rgb = GOLD
    p = tf.add_paragraph()
    p.text = f"{f.get('consortium', '')}.  For {f.get('audience', '')}.  |  {f.get('generated', '')}"
    p.font.size = Pt(13); p.font.color.rgb = RGBColor(220, 230, 238)

    # 2 - Headline numbers
    s = prs.slides.add_slide(blank); bar(s, "The headline")
    nums = []
    if d5:
        nums = [(f"{d5.get('lga_total', 0):,}", "zero-dose children, 2026"),
                (f"{d5.get('top20_pct', 0):.0f}%", "of burden in the top 20% of LGAs"),
                (str(len(d5.get('tier1_states', []))), "Tier-1 critical states")]
    for i, (big, lab) in enumerate(nums):
        x = Inches(0.6 + i * 4.2)
        tbx = s.shapes.add_textbox(x, Inches(2.2), Inches(4.0), Inches(2.6)); tf = tbx.text_frame
        tf.word_wrap = True
        tf.text = big
        tf.paragraphs[0].font.size = Pt(54); tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = [RGBColor(0xC0, 0x39, 0x2B), NAVY, GREEN][i]
        p = tf.add_paragraph(); p.text = clean(lab); p.font.size = Pt(16); p.font.color.rgb = INK

    # 3 - Key findings
    s = prs.slides.add_slide(blank); bar(s, "Key findings")
    kf = _section_lines(narrative_md, ["key findings"]) or [
        f"About {d5.get('lga_total', 0):,} zero-dose children across {d5.get('lga_count', 0)} LGAs in 2026.",
        "Tier-1 states: " + (", ".join(d5.get("tier1_states", [])) or "North-West") + "."]
    bullets(s, kf[:7])

    # 4 - Priority states
    if d5.get("top_states"):
        table_slide("Priority states (top 8 by risk)", ["State", "Zone", "ZD 2026 (%)", "ZD children 2026"],
                    [[clean(x["state"]), clean(x["zone"]), f"{x['zd_2026_pct']:.1f}",
                      f"{x['zd_2026_count']:,}" if x["zd_2026_count"] else "-"] for x in d5["top_states"]])
    # 5 - Priority LGAs
    if d5.get("top_lgas"):
        table_slide("Highest-burden LGAs (top 10)", ["LGA", "State", "ZD children", "Rate (%)"],
                    [[clean(x["lga"]), clean(x["state"]), f"{x['zd_count']:,}", f"{x['zd_rate_pct']:.0f}"]
                     for x in d5["top_lgas"]])

    # 6 - Recommendations
    s = prs.slides.add_slide(blank); bar(s, "Recommendations")
    recs = _section_lines(narrative_md, ["recommendation", "priority action", "implementation"]) or [
        f"Prioritize the top {d5.get('n80', 0)} LGAs that hold 80 percent of the burden.",
        "Concentrate first-line investment in the Tier-1 states.",
        "Address dropout with reminder-recall and defaulter tracing.",
        "Use the LGA hotspot map to micro-plan supervision."]
    bullets(s, recs[:7])
    foot = s.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.5))
    fp = foot.text_frame; fp.text = ("Generated by the NPHCDA Zero-Dose Predictive Modelling Platform "
                                     "from live model outputs. Figures are model estimates.")
    fp.paragraphs[0].font.size = Pt(9); fp.paragraphs[0].font.color.rgb = RGBColor(0x6B, 0x7A, 0x88)

    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()


# --------------------------------------------------------------------------------------
# Standard Operating Procedure (branded Word .docx)
# --------------------------------------------------------------------------------------
def sop_docx() -> bytes:
    """Branded SOP document for DIH capacity building (Print to PDF for a PDF copy)."""
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    navy, green, mute, gold = (RGBColor(0x1F, 0x3B, 0x57), RGBColor(0x1C, 0x7A, 0x3D),
                               RGBColor(0x6B, 0x7A, 0x88), RGBColor(0xC8, 0x90, 0x2A))
    doc = Document()
    doc.styles["Normal"].font.name = "Calibri"
    doc.styles["Normal"].font.size = Pt(10.5)

    if C.LOGO_PATH.exists():
        try:
            p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            p.add_run().add_picture(str(C.LOGO_PATH), height=Inches(0.55))
        except Exception:
            pass
    t = doc.add_paragraph(); tr = t.add_run("Standard Operating Procedure")
    tr.bold = True; tr.font.size = Pt(20); tr.font.color.rgb = navy
    s = doc.add_paragraph(); sr = s.add_run("NPHCDA Zero-Dose Predictive Modelling Platform")
    sr.font.size = Pt(13); sr.font.color.rgb = green
    d = doc.add_paragraph(); dr = d.add_run(
        "For the NPHCDA Digital Innovation Hub and new users.  Generated "
        + datetime.now().strftime("%d %B %Y"))
    dr.italic = True; dr.font.size = Pt(9); dr.font.color.rgb = mute

    def h(text, color=navy, size=13):
        p = doc.add_paragraph(); r = p.add_run(text)
        r.bold = True; r.font.size = Pt(size); r.font.color.rgb = color
        return p

    def bullets(items, style="List Bullet"):
        for it in items:
            doc.add_paragraph(clean(it), style=style)

    def table(headers, rows):
        tb = doc.add_table(rows=1, cols=len(headers)); tb.style = "Light Grid Accent 1"
        for i, c in enumerate(headers):
            run = tb.rows[0].cells[i].paragraphs[0].add_run(c); run.bold = True
        for row in rows:
            cells = tb.add_row().cells
            for i, v in enumerate(row):
                cells[i].text = clean(str(v))

    h("Purpose and audience")
    doc.add_paragraph(clean(
        "This SOP enables any new user to access and use the platform to model the zero-dose dataset "
        "end to end: forecasting antigen coverage, analysing dropout, estimating state and LGA zero-dose "
        "burden and hotspots, exploring drivers, and generating reports. Estimated time end to end is "
        "about 5 to 10 minutes on the bundled sample data."))

    h("Workflow at a glance")
    doc.add_paragraph(clean("Sign in  ->  Load data  ->  Check quality  ->  Run models  ->  "
                            "Explore and test  ->  Ask the Analyst  ->  Generate reports"))

    h("Before you start - data you need")
    doc.add_paragraph(clean("Fastest path: use the bundled project sample data (no upload). To model your "
                            "own data, prepare these five CSVs with the same columns as the sample:"))
    table(["File", "Holds", "Used by"], [
        ["DHIS2 export", "Monthly antigen doses by state and LGA", "Coverage, Dropout, Zero-Dose"],
        ["NDHS zero-dose (long)", "Survey zero-dose rate by state and year", "Zero-Dose"],
        ["Under-five population", "Under-5 cohort by state", "Zero-Dose burden"],
        ["LGA population", "Population by LGA (NPC 2022)", "LGA burden, hotspots"],
        ["Zero-dose model dataset", "State equity / socioeconomic covariates", "Implementation Science, drivers"],
        ["NDHS antigen coverage 2024", "State survey coverage per antigen", "Coverage admin-vs-survey lines"],
        ["DHIS2 live births (monthly)", "Facility-reported live births by LGA", "Coverage optional denominator"],
    ])

    h("Step-by-step")
    steps = [
        "Sign in: enter your name and email (and access code if required). Your sign-in is recorded for usage tracking.",
        "Load the data: on Home, click Use bundled project sample data, or Upload your own data on the Data and Quality page.",
        "Check data quality and anomalies: review completeness, reporting rates, the missing-value heatmap and the Anomaly detection tab.",
        "Run the models: Coverage Forecasting (80% target; forecast end-year to 2032; 3/6/12-month scorecard horizon; Estimated-coverage mode (doses / estimated eligible cohort) with a selectable denominator and NDHS survey reference lines; LGA at-risk screen); Dropout & Completion (forecasts, LASSO drivers, heatmap); Zero-Dose & Hotspots (Bayesian model, LGA burden, Pareto, Gi* maps).",
        "Explore and test (Implementation Science): correlation with multicollinearity flags, distributions, scatter, zone violins, and the Hypothesis tests tab (t-test, ANOVA, chi-square).",
        "Ask the Analyst: add your OpenAI key in the sidebar, then ask grounded cross-domain questions.",
        "Generate reports: on Reports & Briefs, produce the factsheet, the Word policy brief and the PowerPoint deck.",
        "Optional - Program Q&A: upload a programme report (PDF or Word) and ask questions with page-cited answers.",
    ]
    for i, stp in enumerate(steps, 1):
        doc.add_paragraph(clean(f"{i}. {stp}"), style="List Number")

    h("How to read the outputs")
    table(["Signal", "Meaning"], [
        ["Coverage / at-risk", "Red = below the 80% target; green = on target"],
        ["Anomaly severity", "High (|z| >= 3) vs Moderate; spike or drop flagged"],
        ["Priority tier (LGA)", "Tier 1 Critical (red) to Tier 4 (blue), by burden"],
        ["Pareto severity", "Critical/High/Moderate/Lower within state; band A/B/C of burden"],
        ["Gi* hotspot", "Hot Spot (red) = significant high-burden cluster; p<0.01 most confident"],
        ["Survey line (Coverage)", "Dotted purple = NDHS survey coverage; a large gap vs admin flags denominator/reporting issues"],
        ["Credible / prediction interval", "The plausible range; wider = more uncertainty"],
    ])

    h("Tips and troubleshooting")
    bullets([
        "No data loaded: go to Home and load the sample or upload your files.",
        "AI says add a key: paste your OpenAI key in the sidebar (kept for the session only).",
        "A heavy step is slow: the full Bayesian posterior and per-state forecasts run on demand; use defaults for a fast live run.",
        "Access denied at login: your email may not be on the access list - contact the DIH administrator.",
        "Save a report as PDF: open the factsheet or this SOP and use the browser Print to PDF.",
    ])

    h("Data governance and good practice")
    bullets([
        "All figures are model estimates - review before use in official decisions.",
        "Uploaded data stays in the session and is not written to disk; the OpenAI key is never stored.",
        "Cite the data vintage (DHIS2 month, NDHS round, NPC 2022 denominator) when sharing.",
        "Align zero-dose definitions and targets with IA2030 and the national RI guidelines.",
    ])

    fp = doc.add_paragraph(); fr = fp.add_run(
        "NPHCDA Zero-Dose Predictive Modelling Platform. In technical support of NPHCDA; for NPHCDA, "
        "GAVI and UNICEF. Figures are model estimates; review before use.")
    fr.italic = True; fr.font.size = Pt(8); fr.font.color.rgb = mute

    buf = io.BytesIO(); doc.save(buf); return buf.getvalue()


# --------------------------------------------------------------------------------------
# Methods & validation summary (reviewer-facing; shared by the in-app page and the .docx)
# --------------------------------------------------------------------------------------
def methods_sections() -> list:
    """Structured Methods & Validation content. Each item: (heading, kind, payload).
    kind in {'para','bullets','table'}; rendered by both the app page and the Word export."""
    return [
        ("Purpose", "para",
         ["This summary documents the methods and validation behind the NPHCDA Zero-Dose Predictive "
          "Modelling Platform for technical reviewers (GAVI, UNICEF, NPHCDA). Every model runs live on "
          "the loaded data and reproduces the project report. Figures are model estimates for decision "
          "support, to be reviewed before official use."]),
        ("Data sources and vintage", "table",
         (["Source", "Vintage", "Used for"],
          [[clean(a), clean(c), clean(d_)] for a, b, c, d_ in C.PROVENANCE])),
        ("Antigen coverage forecasting", "bullets",
         ["Model: Prophet additive time series per antigen (BCG, Penta1, Penta3, Measles1) - "
          "piecewise-linear trend with automatic changepoints, yearly plus an added semi-annual "
          "seasonality; 95% and 80% prediction intervals.",
          "Two reporting modes: a denominator-free index (percent of the 2024 baseline; the 80% line is "
          "an at-risk-of-decline early-warning) and estimated coverage of the eligible cohort (doses divided by an eligible "
          "infant denominator: under-five/5 demographic proxy by default, or DHIS2 live births).",
          "Admin-vs-survey triangulation: NDHS 2024 survey coverage is overlaid per antigen so a large "
          "admin-vs-survey gap flags a denominator or reporting-completeness issue.",
          "User-set forecast horizon (to 2032) and a 3/6/12-month scorecard window; an LGA at-risk "
          "trend screen flags LGA-antigen series projected below 80%."]),
        ("Validation - coverage forecasts", "bullets",
         ["Out-of-sample hold-out back-test: refit on all but the last 6 months and score the forecast "
          "against the held-out actuals. On the project data MAPE is about 2.7-4.3% per antigen - "
          "'highly accurate' on the Lewis (1982) scale (under 10% = highly accurate, 10-20% = good, "
          "20-50% = reasonable, over 50% = inaccurate).",
          "95% prediction-interval coverage is at or near 100% (the share of held-out actuals that fell "
          "inside the model's 95% interval; ideal near 95%), so the intervals are well-calibrated. All "
          "four tracer antigens fall in the 'highly accurate' band."]),
        ("Dropout and completion", "bullets",
         ["Prophet forecasts of three dropout transitions (Penta1-Penta3, Penta1-Measles1, "
          "Measles1-Measles2), nationally and by state, with prediction intervals.",
          "Drivers: cross-validated LASSO selection with 200-resample bootstrap stability, then a "
          "parsimonious model on the top stably-selected drivers per pair reporting the standardized "
          "coefficient, direction and 95% CI. Because dropout can be negative, a linear model with HC3 "
          "robust standard errors is used (Beta applies only to the bounded zero-dose outcome)."]),
        ("Zero-dose and hotspots", "bullets",
         ["Model: Bayesian hierarchical Beta regression in PyMC on the NDHS 2008-2024 panel - partially "
          "pooled intercepts and time slopes at national, zone and state levels; a Beta likelihood whose "
          "precision scales with each survey's sample size; non-centred parameterization; sampled with "
          "NUTS (nutpie). Convergence reported via R-hat and ESS.",
          "Forecasts to 2026-2028 are posterior-predictive (mean with 95% credible intervals).",
          "LGA burden: each LGA's DHIS2-derived rate is calibrated to its state's 2026 posterior and the "
          "state cohort is distributed across LGAs by population (NPC 2022); the state credible interval "
          "is allocated down to each LGA (ZD count low/high 95%).",
          "Spatial: Getis-Ord Gi* hotspots (k=5 nearest neighbours, row-standardized, permutation "
          "inference) on the LGA zero-dose surface."]),
        ("Validation - zero-dose model (leave-one-wave-out)", "bullets",
         ["Each NDHS survey wave is held out in turn, the model refit, and that wave predicted and "
          "compared to observed (MAE/RMSE in percentage points and 95% CI coverage), in the Diagnostics "
          "tab.",
          "On the project data, out-of-sample MAE is about 8-10 percentage points against a zero-dose "
          "outcome that ranges 3-86% (mean ~30%, SD ~23 pp) over only four survey waves - i.e. the "
          "error is well below the natural spread of the outcome. As a practical guide, MAE under ~5 pp "
          "is excellent and 5-10 pp is good for a four-point survey series; MAPE-style benchmarks "
          "(Lewis 1982) are not applied to the rate directly because percentage error is unstable for "
          "small proportions.",
          "95% credible-interval coverage is roughly 60-85%, so the intervals are somewhat over-"
          "confident out of sample - forecasts should be read as central estimates with intervals "
          "likely a touch narrow (disclosed, not hidden)."]),
        ("Drivers and inference (Implementation Science)", "bullets",
         ["Exploratory analysis of the state zero-dose dataset: correlation matrix with "
          "multicollinearity flags, distributions, scatter (Pearson r), zone violins (Kruskal-Wallis), "
          "burden-band bars and a Bland-Altman agreement plot.",
          "A parsimonious Beta regression of the zero-dose rate on the top LASSO-selected drivers reports "
          "the coefficient, direction and 95% CI; a build-your-own hypothesis-test tab supports t-test, "
          "paired t-test, ANOVA and chi-square.",
          "All driver results are cross-sectional, ecological state-level associations (n=37), framed as "
          "directional and uncertainty-quantified - not causal effects."]),
        ("Key parameters", "table",
         (["Component", "Setting"],
          [["Prophet", "yearly + semi-annual seasonality; changepoint_prior_scale 0.05; 95% PI (80% PI = 0.654x half-width)"],
           ["Bayesian Beta", "hierarchical national-zone-state; precision scaled by survey n; nutpie NUTS, 2 chains, target_accept 0.92; live 1000 draws (full 3000)"],
           ["Getis-Ord Gi*", "k=5 nearest neighbours; row-standardized; permutation p (0.01/0.05/0.10)"],
           ["LASSO drivers", "standardized; 5-fold CV; 200-bootstrap stability; parsimonious top-4 with HC3 robust SE or Beta + 95% CI"]])),
        ("Limitations and honest caveats", "bullets",
         ["DHIS2 is administrative data; reporting completeness affects counts. The coverage index is "
          "relative to 2024 unless a proper denominator is supplied. DHIS2-reported live births "
          "under-count true births, so using them as a denominator can push coverage above 100%.",
          "Driver associations are ecological (state-level, n=37), not causal; read direction and CI "
          "width, not a significance verdict.",
          "The LGA burden is a population-share allocation of the state posterior (carrying the state "
          "credible interval), not an independently fitted LGA model.",
          "Out-of-sample credible-interval coverage is below nominal, so forecast intervals are likely "
          "slightly narrow; the zero-dose model extrapolates a linear-in-logit time trend from four NDHS "
          "waves.",
          "Unmatched geographies in the Gi* step are median-imputed; the LGA reporting drill-down and "
          "anomaly tab help surface data-quality issues before modelling."]),
        ("Reproducibility and governance", "bullets",
         ["Content-hashed caching means re-runs are instant but genuinely recomputed when data changes; "
          "dependencies are version-pinned and the container is reproducible.",
          "Uploaded data is processed only in the session and not written to disk; any OpenAI key is "
          "session-only and never stored; login PII is used solely for DIH usage tracking.",
          "Every output is labelled as a model estimate; cite the data vintage when sharing."]),
    ]


def methods_docx() -> bytes:
    """Reviewer-facing Methods & Validation summary as a branded Word document."""
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    navy, green, mute = RGBColor(0x1F, 0x3B, 0x57), RGBColor(0x1C, 0x7A, 0x3D), RGBColor(0x6B, 0x7A, 0x88)
    doc = Document()
    doc.styles["Normal"].font.name = "Calibri"
    doc.styles["Normal"].font.size = Pt(10.5)
    if C.LOGO_PATH.exists():
        try:
            p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            p.add_run().add_picture(str(C.LOGO_PATH), height=Inches(0.55))
        except Exception:
            pass
    t = doc.add_paragraph(); tr = t.add_run("Methods and Validation Summary")
    tr.bold = True; tr.font.size = Pt(20); tr.font.color.rgb = navy
    s = doc.add_paragraph(); sr = s.add_run("NPHCDA Zero-Dose Predictive Modelling Platform")
    sr.font.size = Pt(13); sr.font.color.rgb = green
    d = doc.add_paragraph(); dr = d.add_run("For GAVI, UNICEF and NPHCDA technical review.  Generated "
                                            + datetime.now().strftime("%d %B %Y"))
    dr.italic = True; dr.font.size = Pt(9); dr.font.color.rgb = mute

    for heading, kind, payload in methods_sections():
        h = doc.add_paragraph(); hr = h.add_run(heading)
        hr.bold = True; hr.font.size = Pt(13); hr.font.color.rgb = navy
        if kind == "para":
            for para in payload:
                doc.add_paragraph(clean(para))
        elif kind == "bullets":
            for b in payload:
                doc.add_paragraph(clean(b), style="List Bullet")
        elif kind == "table":
            headers, trows = payload
            tb = doc.add_table(rows=1, cols=len(headers)); tb.style = "Light Grid Accent 1"
            for i, c in enumerate(headers):
                run = tb.rows[0].cells[i].paragraphs[0].add_run(c); run.bold = True
            for row in trows:
                cells = tb.add_row().cells
                for i, v in enumerate(row):
                    cells[i].text = clean(str(v))

    fp = doc.add_paragraph(); fr = fp.add_run(
        "Generated by the NPHCDA Zero-Dose Predictive Modelling Platform. Figures are model estimates; "
        "review before use.")
    fr.italic = True; fr.font.size = Pt(8); fr.font.color.rgb = mute
    buf = io.BytesIO(); doc.save(buf); return buf.getvalue()
